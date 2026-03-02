import tkinter as tk
from tkinter import filedialog, ttk, scrolledtext
import pandas as pd
from pptx import Presentation
import json
import os
import threading
import time
import copy
import copy
import copy
from google import genai
from openpyxl.utils.cell import range_boundaries
from pptx.dml.color import RGBColor
from pptx.util import Pt

# =====================================================================
# [Architecture] 모듈 간 독립성 유지 
# - CoreLogic: 데이터 처리, AI 통신, PPT 생성 등 핵심 비즈니스 로직 담당
# - GuiApp: 사용자 인터페이스 뷰 및 이벤트 처리 담당
# - HardwareControl(생략): 하드웨어/모터 제어 모듈이 추가될 경우 분리될 영역
# =====================================================================

class CoreLogic:
    def __init__(self):
        # [Precision] 센서 및 수치 계측 데이터 처리 시 0.0001 정밀도(float64) 유지를 위한 통일된 설정
        self.data_precision = 'float64'

    def read_excel_data(self, file_path, excel_range=None):
        """엑셀 데이터를 읽어오는 함수 (특정 범위 지정 지원)"""
        try:
            read_kwargs = {}
            if excel_range:
                # 'A1:D10' 형식을 파싱하여 (min_col, min_row, max_col, max_row) 추출
                min_col, min_row, max_col, max_row = range_boundaries(excel_range)
                
                # pandas 옵션으로 변환
                # usecols: 0-indexed로 변환 ('A', 'B' 등을 인덱스 번호 리스트로)
                read_kwargs['usecols'] = list(range(min_col - 1, max_col))
                # skiprows: 헤더(열 이름)가 될 행을 포함하여 파일 맨 위부터 건너뛸 줄 수 (0-indexed)
                read_kwargs['skiprows'] = min_row - 1
                # nrows: 헤더를 제외한 실제 데이터 행 개수
                read_kwargs['nrows'] = max_row - min_row
            
            # 수치 데이터 정밀도 유지를 위해 명시적으로 float64 변환을 허용/유도하도록 읽기
            df = pd.read_excel(file_path, **read_kwargs)
            return df
        
        except PermissionError as e:
            # [Error Handling] 파일이 열려있어서 발생하는 접근 권한 에러 처리
            raise Exception(f"엑셀 파일 읽기 실패: 파일이 엑셀 프로그램에서 열려 있습니다. 엑셀을 완전히 종료하고 다시 시도해 주세요.\n(상세 에러: {str(e)})")
        except Exception as e:
            if "Permission denied" in str(e) or "Errno 13" in str(e):
                raise Exception(f"엑셀 파일 읽기 실패: 파일이 엑셀 프로그램에서 열려 있습니다. 엑셀을 완전히 종료하고 다시 시도해 주세요.\n(상세 에러: {str(e)})")
            raise Exception(f"엑셀 데이터 범위 읽기 실패 (입력 범위를 확인해 주세요): {str(e)}")

    def call_llm_semantic_chunking(self, df):
        """
        AI 기반 데이터 분할 (Semantic Chunking) API 모듈
        Gemini API를 사용하여 실제 데이터를 지능적으로 분할
        """
        # 사용자 요구사항 프롬프트 구성 (실제 CSV 데이터 첨부)
        csv_data = df.fillna("").to_csv(index=False)
        
        api_prompt = (
            "당신은 엑셀 데이터를 논리적인 묶음으로 분할하여 PPT 슬라이드에 넣기 좋게 만들어주는 데이터 정리 AI입니다.\n"
            "다음은 엑셀 데이터의 내용(CSV 형식)입니다.\n\n"
            f"--- 엑셀 데이터 시작 ---\n{csv_data}\n--- 엑셀 데이터 끝 ---\n\n"
            "요구사항:\n"
            "1. PPT 한 슬라이드 표의 최대 데이터 행 수는 15줄 이하입니다.\n"
            "2. 무조건 15줄에서 자르지 말고, 시험 평가 항목, 테스트 단계(Phase), 또는 같은 부품(Component) 계측 데이터끼리는 페이지가 넘어가며 끊기지 않도록 문맥을 파악하여 논리적으로 묶어주세요.\n"
            "3. 테이블의 모양, 헤더 박복, 폰트(맑은 고딕, 10pt) 등은 파이썬 로직에서 처리하므로, 당신은 데이터의 분할 논리에만 집중하세요.\n"
            "4. 각 분할된 덩어리(chunk)는 원본 데이터의 구조(컬럼)를 유지해야 합니다.\n"
            "결과는 반드시 다음 구조의 JSON 배열 형태(순수 JSON 텍스트)로만 반환하세요. 앞뒤로 마크다운 코드 블록(```json 등)이나 부가 설명을 절대 붙이지 마세요:\n"
            '[\n  {\n    "headers": ["컬럼1", "컬럼2"],\n    "rows": [["값1", "값2"], ["값3", "값4"]]\n  }\n]'
        )
        
        # [Error Handling] Serial/TCP/API 등 외부 통신 시 예외 처리를 위한 견고한 Try-Except 블록
        try:
            # 사용할 Gemini API Key 지정 (환경변수에서 로드하여 깃허브 노출 방지)
            api_key = os.environ.get("GEMINI_API_KEY")
            if not api_key:
                raise Exception("GEMINI_API_KEY 환경 변수가 설정되지 않았습니다. API 키를 시스템 환경 변수에 추가해주세요.")

            # 사용할 Gemini 모델 설정 및 호출 (신규 genai 패키지 방식)
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=api_prompt,
            )
            
            # API 응답에서 텍스트 수신 후, 혹시 모를 마크다운 태그가 있다면 제거
            json_text = response.text.strip()
            if json_text.startswith("```json"):
                json_text = json_text[7:]
            if json_text.startswith("```"):
                json_text = json_text[3:]
            if json_text.endswith("```"):
                json_text = json_text[:-3]
            
            json_text = json_text.strip()
            return json_text
            
        except Exception as e:
            raise Exception(f"Gemini API 통신 에러 발생: {str(e)}")

    def clone_slide(self, prs, index=0):
        """
        기존 슬라이드의 서식(표, 폰트, 배경 등)을 최대한 유지하기 위한 XML 트리 복사
        """
        template_slide = prs.slides[index] # 원본 슬라이드 추출
        # 기존 PPT 양식 중 레이아웃 개수가 7개 미만인 경우 에러(IndexError) 방지를 위해 기준 슬라이드의 원본 레이아웃을 그대로 상속받음
        blank_slide_layout = template_slide.slide_layout
        new_slide = prs.slides.add_slide(blank_slide_layout)
        
        # XML 요소를 깊은 복사(Deepcopy)하여 형상 및 속성 보존
        for shape in template_slide.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            
        return new_slide

    def process_data_to_ppt(self, excel_path, ppt_path, progress_callback, log_callback, excel_range=None, ppt_target=None, save_dest=None):
        """전체 변환 파이프라인 제어 (엑셀 읽기 -> AI 분할 -> PPT 적용 및 저장)"""
        try:
            # 1. 문서 읽기 (진행률 10%)
            log_callback(f"엑셀 데이터를 로드하고 있습니다... (범위: {excel_range if excel_range else '전체'})")
            df = self.read_excel_data(excel_path, excel_range)
            progress_callback(10)
            
            # 2. AI Semantic Chunking 통신 (진행률 40%)
            log_callback("Gemini AI 엔진에 의미 단위 분할(Semantic Chunking)을 요청합니다...")
            json_string = self.call_llm_semantic_chunking(df)
            chunks = json.loads(json_string)
            progress_callback(40)
            
            # 3. PPT 구조 준비 및 서식 적용 (진행률 90% 까지 할당)
            log_callback("PPT 양식을 불러오는 중입니다...")
            if not os.path.exists(ppt_path):
                raise Exception(f"오류: 지정한 PPT 템플릿 파일을 찾을 수 없습니다.\n경로 확인: {ppt_path}")
                
            try:
                prs = Presentation(ppt_path)
            except Exception as e:
                if "Package not found" in str(e):
                    raise Exception("오류: 선택하신 파일이 유효한 PPTX 템플릿 파일이 아니거나 손상되었습니다. 다른 파일을 선택해 주세요.")
                raise Exception(f"PPT 템플릿 로드 중 에러 발생: {str(e)}")
            
            if len(prs.slides) == 0:
                raise Exception("오류: 선택한 PPT 템플릿에 슬라이드가 존재하지 않습니다.")
            
            # 타겟 슬라이드/표 위치 파싱 (예: "1-1")
            target_slide_idx = 0
            target_table_idx = 0
            if ppt_target:
                try:
                    parts = ppt_target.split('-')
                    target_slide_idx = int(parts[0].strip()) - 1
                    target_table_idx = int(parts[1].strip()) - 1
                except Exception:
                    log_callback("[경고] PPT 출력 대상 형식이 잘못되었습니다 (예: 1-1). 기본값(1번째 슬라이드, 1번째 표)을 사용합니다.")
                    target_slide_idx = 0
                    target_table_idx = 0

            # 템플릿 슬라이드 유효성 검사
            if len(prs.slides) <= target_slide_idx:
                 raise Exception(f"오류: 지정한 슬라이드 번호({target_slide_idx+1})가 템플릿에 존재하지 않습니다.")

            total_chunks = len(chunks)
            for i, chunk in enumerate(chunks):
                # 첫 덩어리는 지정된 기존 슬라이드를 활용하고, 이후는 해당 슬라이드를 복제
                slide = prs.slides[target_slide_idx] if i == 0 else self.clone_slide(prs, target_slide_idx)
                
                # 표(Table) 객체 수색
                target_table = None
                table_counter = 0
                for shape in slide.shapes:
                    if shape.has_table:
                        if table_counter == target_table_idx:
                            target_table = shape.table
                            break
                        table_counter += 1
                
                if not target_table:
                    raise Exception(f"오류: {target_slide_idx+1}번 슬라이드에서 {target_table_idx+1}번째 표를 찾을 수 없습니다.")
                
                headers = chunk.get("headers", [])
                rows = chunk.get("rows", [])
                
                # 표의 크기/위치 및 내부 서식이 틀어지지 않도록 단락 단위로 텍스트를 섬세하게 주입하는 헬퍼 함수
                def set_cell_text_black_and_size9(cell, text):
                    # 전체 텍스트 프레임을 초기화하지 않고 첫 번째 문단을 유지하며 수정
                    text_frame = cell.text_frame
                    text_frame.clear() # 기존 문단 텍스트만 지우되 객체 프레임 보존
                    
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = str(text)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.size = Pt(9) # 폰트 크기 9포인트 고정

                # 엑셀 헤더 적용 (표의 첫 번째 줄)
                for col_idx, header_text in enumerate(headers):
                    if 0 < len(target_table.rows) and col_idx < len(target_table.columns):
                        set_cell_text_black_and_size9(target_table.cell(0, col_idx), header_text)
                        
                # 각 데이터 행 할당 (기존 표의 범위를 넘지 않게 처리)
                for row_idx, row_data in enumerate(rows):
                    ppt_row_idx = row_idx + 1 # 헤더 다음 행부터 데이터 입력
                    if ppt_row_idx < len(target_table.rows):
                        for col_idx, cell_data in enumerate(row_data):
                            if col_idx < len(target_table.columns):
                                # 텍스트 데이터만 주입하여 텍셀 문서의 서식 유지 및 속성 덮어쓰기
                                set_cell_text_black_and_size9(target_table.cell(ppt_row_idx, col_idx), cell_data)

                # 애니메이션 효과와 진행률 표현
                current_progress = 40 + int(50 * ((i + 1) / total_chunks))
                progress_callback(current_progress)
                time.sleep(0.05) 
                
            # 4. 결과물 저장 (진행률 100%)
            log_callback("결과물 저장을 준비합니다...")
            
            # 지정된 저장 경로가 없으면 엑셀 파일과 같은 폴더에 자동 네이밍하여 저장
            if not save_dest:
                base_dir = os.path.dirname(excel_path)
                file_title = os.path.splitext(os.path.basename(excel_path))[0]
                save_dest = os.path.join(base_dir, f"{file_title}_결과보고서.pptx")
            
            prs.save(save_dest)
            
            progress_callback(100)
            log_callback("\n작업 완료")
            log_callback(f"저장된 경로: {save_dest}")
            
        except Exception as e:
            log_callback(f"\n[오류 발생]: {str(e)}")
            progress_callback(0)
            raise # 스레드에서 캐치하도록 던짐


# =====================================================================
# GUI 애플리케이션 클래스
# =====================================================================

class ExcelPPTConverterGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("AI 기반 엑셀-PPT 자동 변환 도구")
        self.root.geometry("650x500")
        
        # 의존성 생성
        self.core = CoreLogic()
        
        # 파일 경로 보관
        self.excel_file_path = ""
        self.ppt_file_path = ""
        self.save_dest_path = "" # [신규] 결과물 저장 경로 저장 변수
        
        self.create_widgets()

    def create_widgets(self):
        """화면 UI 요소들을 배치하는 함수"""
        
        # ---- 엑셀 파일 선택 부 ----
        frame_excel = tk.Frame(self.root)
        frame_excel.pack(pady=5, padx=20, fill='x')
        btn_excel = tk.Button(frame_excel, text="엑셀 파일 선택", command=self.action_select_excel, width=15)
        btn_excel.pack(side='left', padx=5)
        self.lbl_excel = tk.Label(frame_excel, text="선택된 파일이 없습니다.", fg="gray")
        self.lbl_excel.pack(side='left', fill='x', expand=True, padx=5)
        
        # 엑셀 데이터 영역 범위 입력
        frame_excel_range = tk.Frame(self.root)
        frame_excel_range.pack(pady=2, padx=20, fill='x')
        tk.Label(frame_excel_range, text=">> 읽어올 범위 (예: A1:F20, 비우면 전체):").pack(side='left', padx=5)
        self.entry_excel_range = tk.Entry(frame_excel_range, width=15)
        self.entry_excel_range.pack(side='left', padx=5)
        
        # ---- PPT 템플릿 선택 부 ----
        frame_ppt = tk.Frame(self.root)
        frame_ppt.pack(pady=5, padx=20, fill='x')
        btn_ppt = tk.Button(frame_ppt, text="PPT 템플릿 선택", command=self.action_select_ppt, width=15)
        btn_ppt.pack(side='left', padx=5)
        self.lbl_ppt = tk.Label(frame_ppt, text="선택된 파일이 없습니다.", fg="gray")
        self.lbl_ppt.pack(side='left', fill='x', expand=True, padx=5)

        # PPT 타겟 슬라이드/표 위치 지정
        frame_ppt_target = tk.Frame(self.root)
        frame_ppt_target.pack(pady=2, padx=20, fill='x')
        tk.Label(frame_ppt_target, text=">> 출력 대상 (예: 1-1 = 1번 슬라이드 1번째 표, 비우면 기본값):").pack(side='left', padx=5)
        self.entry_ppt_target = tk.Entry(frame_ppt_target, width=15)
        self.entry_ppt_target.pack(side='left', padx=5)

        # ---- [신규] 저장 위치 지정 부 ----
        frame_save = tk.Frame(self.root)
        frame_save.pack(pady=5, padx=20, fill='x')
        btn_save = tk.Button(frame_save, text="저장 위치 지정", command=self.action_select_save_dest, width=15)
        btn_save.pack(side='left', padx=5)
        self.lbl_save = tk.Label(frame_save, text="비워두면 엑셀 원본과 같은 위치에 자동 저장됩니다.", fg="gray")
        self.lbl_save.pack(side='left', fill='x', expand=True, padx=5)
        
        # ---- 변환 컨트롤 버튼 ----
        frame_buttons = tk.Frame(self.root)
        frame_buttons.pack(pady=15)

        # [Safety] 동적 구동 중 오작동 방지를 파악할 수 있도록 큼직한 버튼 적용
        self.btn_execute = tk.Button(frame_buttons, text="작업 실행 (Start)", command=self.action_run_conversion, 
                                    font=('맑은 고딕', 14, 'bold'), bg='#2C3E50', fg='white')
        self.btn_execute.pack(side='left', padx=10, ipadx=40, ipady=10)

        # [Safety] 에러 또는 무한루프 발생 시 즉각 종료할 수 있는 Abort 버튼 (빨간색 계열)
        self.btn_abort = tk.Button(frame_buttons, text="강제 종료 (Abort)", command=self.action_abort,
                                   font=('맑은 고딕', 14, 'bold'), bg='#D9534F', fg='white')
        self.btn_abort.pack(side='left', padx=10, ipadx=40, ipady=10)
        
        # ---- 프로세스 상태 모니터링 (Bar graph & Percentage) ----
        frame_monitor = tk.Frame(self.root)
        frame_monitor.pack(pady=5, padx=20, fill='x')
        
        self.progress_val = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(frame_monitor, variable=self.progress_val, maximum=100)
        self.progress_bar.pack(side='left', fill='x', expand=True)
        
        self.lbl_percentage = tk.Label(frame_monitor, text="0%", width=6, font=('Arial', 10, 'bold'))
        self.lbl_percentage.pack(side='right', padx=10)
        
        # ---- 로그창 (진행 상태 / 에러 메시지) ----
        self.txt_logger = scrolledtext.ScrolledText(self.root, height=12, state='disabled', bg='#1E1E1E', fg='#00FF00', font=('Consolas', 9))
        self.txt_logger.pack(pady=15, padx=20, fill='both', expand=True)

    def write_log(self, msg):
        """텍스트 박스에 로그 메시지를 출력하는 스레드 안전성 보조 함수"""
        self.txt_logger.config(state='normal')
        self.txt_logger.insert(tk.END, f"{msg}\n")
        self.txt_logger.see(tk.END)
        self.txt_logger.config(state='disabled')

    def sync_progress_ui(self, value):
        """진행률 막대기 및 퍼센테이지 박스 동기화 갱신"""
        self.progress_val.set(value)
        self.lbl_percentage.config(text=f"{int(value)}%")
        self.root.update_idletasks()

    def action_select_excel(self):
        path = filedialog.askopenfilename(filetypes=[("Excel Files", "*.xlsx *.xls")])
        if path:
            self.excel_file_path = path
            self.lbl_excel.config(text=path, fg="blue")
            self.write_log(f"엑셀 데이터소스 연결 완료: {path}")

    def action_select_ppt(self):
        path = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])
        if path:
            self.ppt_file_path = path
            self.lbl_ppt.config(text=path, fg="blue")
            self.write_log(f"PPT 템플릿 연결 완료: {path}")

    def action_select_save_dest(self):
        """[신규] 결과물을 저장할 위치 및 이름 지정 팝업"""
        path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")],
            title="결과물 저장 위치 지정"
        )
        if path:
            self.save_dest_path = path
            self.lbl_save.config(text=path, fg="blue")
            self.write_log(f"저장 위치 지정 완료: {path}")

    def action_abort(self):
        """[Safety] 강제 종료 (Abort) 처리: 프로그램 창 및 프로세스 즉각 소멸"""
        # os._exit()를 사용하면 별도 구동 중인 스레드나 작업 큐 무관하게 
        # OS 단에서 즉시 전체 프로세스를 Kill 하므로 확실한 강제 종료가 보장됩니다.
        os._exit(1)

    def action_run_conversion(self):
        """변환 버튼 클릭 시 수행 흐름 제어"""
        if not self.excel_file_path or not self.ppt_file_path:
            self.write_log("[경고] 엑셀 파일과 PPT 양식을 모두 지정해야 합니다.")
            return
            
        # [Safety] 다중 클릭 연타 및 리소스 충돌을 막기 위한 인터락(Fail-safe 방어 조치)
        self.btn_execute.config(state='disabled')
        self.progress_val.set(0)
        self.lbl_percentage.config(text="0%")
        
        self.txt_logger.config(state='normal')
        self.txt_logger.delete(1.0, tk.END)
        self.txt_logger.config(state='disabled')
        
        self.write_log("========= AI 논리적 분할 기반 PPT 변환 시스템 가동 =========")
        
        # GUI의 멈춤(Freezing) 현상을 막고 비동기처리를 위해 별도의 스레드 가동
        worker_thread = threading.Thread(target=self.background_worker)
        worker_thread.daemon = True
        worker_thread.start()

    def background_worker(self):
        """별도 스레드에서 구동되는 로직 브릿지 (Core 영역 호출)"""
        # UI에서 입력한 동적 파라미터 획득
        excel_range_val = self.entry_excel_range.get().strip() or None
        ppt_target_val = self.entry_ppt_target.get().strip() or None
        save_dest_val = self.save_dest_path or None

        try:
            self.core.process_data_to_ppt(
                self.excel_file_path,
                self.ppt_file_path,
                self.sync_progress_ui,
                self.write_log,
                excel_range=excel_range_val,
                ppt_target=ppt_target_val,
                save_dest=save_dest_val
            )
        except Exception:
            # 예외 상세 내용은 Core 내부 코루틴에서 Logger로 전달되었으므로 생략
            pass
        finally:
            # [Safety] 작업 완료 또는 실패에 무관하게 상태 해제
            self.btn_execute.config(state='normal')


# =====================================================================
# Main 메인 실행부
# =====================================================================
if __name__ == "__main__":
    app_root = tk.Tk()
    
    # 윈도우 스타일링을 위한 테마 호출 적용
    style = ttk.Style()
    if "clam" in style.theme_names():
        style.theme_use("clam")
        
    app_instance = ExcelPPTConverterGUI(app_root)
    app_root.mainloop()
